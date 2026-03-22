"""
Document text extractor — supports PDF, PPTX, TXT, and MD files.
Adapted from zoom-agent/backend/extractor.py.
"""

from __future__ import annotations

import io


def extract_text(filename: str, content: bytes) -> str:
    """Extract plain text from a document.

    Supports PDF (via pymupdf), PPTX (via python-pptx), and TXT/MD.
    Returns an empty string on unsupported file types.
    """
    lower = filename.lower()

    if lower.endswith(".pdf"):
        return _extract_pdf(content)

    if lower.endswith(".pptx") or lower.endswith(".ppt"):
        return _extract_pptx(content)

    if lower.endswith(".txt") or lower.endswith(".md"):
        return content.decode("utf-8", errors="replace")

    return ""


def _extract_pdf(content: bytes) -> str:
    try:
        import fitz  # pymupdf

        doc = fitz.open(stream=content, filetype="pdf")
        pages: list[str] = []
        for page in doc:
            pages.append(page.get_text())
        return "\n".join(pages)
    except ImportError:
        raise RuntimeError(
            "PDF extraction requires 'pymupdf'. Install with: pip install pymupdf"
        )
    except Exception as exc:
        raise RuntimeError(f"PDF extraction failed: {exc}") from exc


def _extract_pptx(content: bytes) -> str:
    try:
        from pptx import Presentation

        prs = Presentation(io.BytesIO(content))
        slides: list[str] = []
        for slide in prs.slides:
            texts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = " ".join(run.text for run in para.runs).strip()
                        if line:
                            texts.append(line)
            if texts:
                slides.append("\n".join(texts))
        return "\n\n".join(slides)
    except ImportError:
        raise RuntimeError(
            "PPTX extraction requires 'python-pptx'. Install with: pip install python-pptx"
        )
    except Exception as exc:
        raise RuntimeError(f"PPTX extraction failed: {exc}") from exc
